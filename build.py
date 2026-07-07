# -*- coding: utf-8 -*-
"""제주교육 AI · 공개 문서 사이트 빌드 (자체 완결 HTML + Bloom PDF + 검증).

- 스크립트가 있는 저장소 폴더(ROOT) 기준으로만 동작 → 경로 하드코딩 없음, 어느 PC에서나 실행.
- 각 .md 를 임베디드 CSS HTML 로 변환하고 index.html 생성.
- Bloom 페이지는 Chrome/Edge 헤드리스로 PDF 저장 후 검증(빈 페이지·에러 페이지 차단).

필요: python -m pip install markdown pypdf   /  Chrome 또는 Edge 설치
사용: python build.py
"""
import os
import re
import sys
import subprocess
from pathlib import Path
import markdown
import html as _html
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent

# (소스 md 파일명, 출력 슬러그, 브라우저 탭 제목, PDF 슬러그 또는 None, PPTX 슬러그 또는 None)
PAGES = [
    ("AI교육-패러다임.md", "paradigm.html", "AI 교육 패러다임", None, "paradigm.pptx"),
    ("AX전환과 교육.md", "ax-transformation.html", "AX 전환과 교육", None, "ax-transformation.pptx"),
    ("고의숙 교육감의 초개별화 맞춤형 교육 이론적 배경.md", "bloom-2sigma.html",
     '고의숙 교육감의 "초개별화 맞춤형 교육"', "bloom-2sigma.pdf", "bloom-2sigma.pptx"),
]

# 마크다운 내부 링크(.md) → 사이트 내 .html 슬러그(구 파일명도 함께 매핑해 안전)
LINK_MAP = {
    "AI교육-패러다임.md": "paradigm.html",
    "AX전환과 교육.md": "ax-transformation.html",
    "AX전환-요약.md": "ax-transformation.html",
    "고의숙 교육감의 초개별화 맞춤형 교육 이론적 배경.md": "bloom-2sigma.html",
    "초개별화 맞춤형 교육의 근거 -벤자민 블룸-2시그마-문제-논문-요약.md": "bloom-2sigma.html",
}

CHROME_CANDIDATES = [
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
]

CSS = """
:root{--bg:#ffffff;--fg:#1a1a1a;--muted:#5b6470;--line:#e3e6ea;--accent:#1f6feb;--soft:#f6f8fa;--quote:#f0f4f9}
@media (prefers-color-scheme:dark){:root{--bg:#0d1117;--fg:#e6edf3;--muted:#9aa4b2;--line:#283039;--accent:#58a6ff;--soft:#161b22;--quote:#12181f}}
*{box-sizing:border-box}
html{-webkit-text-size-adjust:100%}
body{margin:0;background:var(--bg);color:var(--fg);
 font-family:"Pretendard Variable","Pretendard",-apple-system,BlinkMacSystemFont,"Segoe UI","Malgun Gothic","Apple SD Gothic Neo","Noto Sans KR",sans-serif;
 line-height:1.72;font-size:17px;word-break:keep-all;overflow-wrap:anywhere}
.wrap{max-width:820px;margin:0 auto;padding:32px 22px 80px}
nav.top{max-width:820px;margin:0 auto;padding:12px 22px;border-bottom:1px solid var(--line);font-size:14px}
nav.top a.home{display:inline-flex;align-items:center;gap:10px;color:var(--muted);text-decoration:none}
nav.top a.home:hover{color:var(--accent)}
nav.top img.logo{height:36px;width:auto;display:block;background:#fff;padding:3px 6px;border-radius:6px;border:1px solid var(--line)}
h1{font-size:1.9rem;line-height:1.35;margin:.2em 0 .5em;letter-spacing:-.01em}
h1[align="center"]{margin:.1em 0 .15em}
.lead-sub{color:var(--muted);font-size:1.05rem;font-weight:600;margin:-.2em 0 1.4em;text-align:center}
p[align="center"] sub{vertical-align:baseline;font-size:1.05rem;color:var(--muted);font-weight:600}
p[align="center"]{margin:0 0 1.4em}
h2{font-size:1.4rem;margin:1.8em 0 .5em;padding-bottom:.25em;border-bottom:1px solid var(--line)}
h3{font-size:1.15rem;margin:1.4em 0 .4em}
p{margin:.7em 0}
a{color:var(--accent);text-decoration:none}
a:hover{text-decoration:underline}
ul,ol{padding-left:1.4em;margin:.6em 0}
li{margin:.25em 0}
blockquote{margin:1em 0;padding:.6em 1em;background:var(--quote);border-left:4px solid var(--accent);border-radius:4px}
blockquote p{margin:.4em 0}
code{background:var(--soft);padding:.12em .4em;border-radius:4px;font-size:.9em;
 font-family:"SFMono-Regular",Consolas,"Liberation Mono",monospace}
pre{background:var(--soft);padding:14px 16px;border-radius:8px;overflow-x:auto}
pre code{background:none;padding:0}
hr{border:none;border-top:1px solid var(--line);margin:2em 0}
.tablewrap{overflow-x:auto;margin:1em 0}
table{border-collapse:collapse;width:100%;font-size:.94em}
th,td{border:1px solid var(--line);padding:8px 11px;text-align:left;vertical-align:top}
th{background:var(--soft);font-weight:600}
tr:nth-child(even) td{background:color-mix(in srgb,var(--soft) 45%,transparent)}
.foot{max-width:820px;margin:0 auto;padding:24px 22px;color:var(--muted);font-size:13px;border-top:1px solid var(--line)}
.pdflink{display:inline-block;margin:0 0 18px;padding:8px 14px;border:1px solid var(--line);
 border-radius:8px;background:var(--soft);color:var(--fg);font-size:14px;font-weight:600}
.pdflink:hover{border-color:var(--accent);color:var(--accent);text-decoration:none}
.idx-pdf{margin-left:8px;font-size:13px;padding:1px 8px;border:1px solid var(--line);border-radius:12px}
@media print{nav.top,.foot,.pdflink{display:none}body{font-size:11pt;font-family:"Malgun Gothic","Apple SD Gothic Neo","Noto Sans KR",sans-serif}.wrap{max-width:none}}
"""

TEMPLATE = """<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<link rel="icon" href="logo-jje.jpg">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/gh/orioncactus/pretendard@v1.3.9/dist/web/variable/pretendardvariable-dynamic-subset.min.css">
<style>{css}</style>
</head>
<body>
<nav class="top"><a class="home" href="index.html" title="제주도교육인수위원회 미래학력분과 홈"><img class="logo" src="logo-jje.jpg" alt="제주특별자치도교육청"><span>제주도교육인수위원회 미래학력분과</span></a></nav>
<main class="wrap">
{body}
</main>
<div class="foot">원문 마크다운에서 자동 생성</div>
</body>
</html>
"""

md = markdown.Markdown(extensions=["tables", "fenced_code", "sane_lists", "attr_list"])


def convert(md_name):
    text = (ROOT / md_name).read_text(encoding="utf-8")
    for k, v in LINK_MAP.items():
        text = text.replace(k, v)
    md.reset()
    html = md.convert(text)
    html = re.sub(r"<table>", '<div class="tablewrap"><table>', html)
    html = re.sub(r"</table>", "</table></div>", html)
    return html


def find_browser():
    for p in CHROME_CANDIDATES:
        if os.path.exists(p):
            return p
    return None


def make_pdf(slug, pdf_name):
    browser = find_browser()
    if not browser:
        sys.exit("Chrome/Edge 를 찾지 못했습니다.")
    html_uri = (ROOT / slug).resolve().as_uri()   # file:///C:/... (경로 인코딩 안전)
    pdf_path = ROOT / pdf_name
    subprocess.run(
        [browser, "--headless=new", "--no-sandbox", "--disable-gpu",
         "--no-pdf-header-footer", f"--print-to-pdf={pdf_path}", html_uri],
        check=True, capture_output=True)
    # 검증: 정상 PDF는 크고 다중 페이지. 에러 페이지(ERR_FILE_NOT_FOUND)는 1쪽·소용량.
    # (웹폰트 임베드 시 한글 텍스트 추출이 깨질 수 있어 '본문 키워드'는 실패 근거로 쓰지 않음)
    data = pdf_path.read_bytes()
    ok = data[:5] == b"%PDF-" and len(data) > 80000
    pages = "?"
    try:
        from pypdf import PdfReader
        r = PdfReader(str(pdf_path))
        pages = len(r.pages)
        txt = "".join((pg.extract_text() or "") for pg in r.pages)
        if "ERR_FILE" in txt or "cannot be reached" in txt.lower():
            ok = False
        if isinstance(pages, int) and pages < 1:
            ok = False
    except ImportError:
        pass
    print(f"  PDF {pdf_name}: {len(data):,}B, pages={pages}, ok={ok}")
    if not ok:
        sys.exit(f"PDF 검증 실패: {pdf_name} (빈 페이지/에러 페이지 의심)")


def _clean(s):
    """마크다운 인라인 서식 제거 → 슬라이드용 순수 텍스트."""
    s = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", s)              # 이미지
    s = re.sub(r"\[([^\]]+)\]\(<?[^>)]*>?\)", r"\1", s)     # 링크 → 텍스트
    s = s.replace("**", "").replace("`", "")
    s = re.sub(r"\*(.+?)\*", r"\1", s)                       # *강조*
    s = re.sub(r"<[^>]+>", "", s)                            # 잔여 HTML 태그
    return _html.unescape(s).strip()


def _add_table_slide(prs, title, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    ncols = max(len(r) for r in rows)
    nrows = len(rows)
    gt = slide.shapes.add_table(nrows, ncols, Inches(0.4), Inches(1.3),
                                Inches(9.2), Inches(0.4 * nrows)).table
    for i, r in enumerate(rows):
        for j in range(ncols):
            cell = gt.cell(i, j)
            cell.text = r[j] if j < len(r) else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                if i == 0:
                    p.font.bold = True


def make_pptx(md_name, pptx_name, deck_title):
    """소스 .md 를 섹션(##) 단위 슬라이드 덱으로 변환."""
    lines = (ROOT / md_name).read_text(encoding="utf-8").split("\n")
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)

    ts = prs.slides.add_slide(prs.slide_layouts[0])
    ts.shapes.title.text = _clean(deck_title)
    try:
        ts.placeholders[1].text = "제주도교육인수위원회 미래학력분과"
    except (KeyError, IndexError):
        pass

    sections, cur, blocks, tbl, in_code = [], None, [], [], False

    def flush_tbl():
        if tbl:
            blocks.append(("table", list(tbl)))
            tbl.clear()

    for raw in lines:
        st = raw.strip()
        if st.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            continue
        if st.startswith("## "):
            flush_tbl()
            if cur is not None:
                sections.append((cur, blocks))
            cur, blocks = _clean(st[3:]), []
            continue
        if st.startswith("### "):
            flush_tbl()
            blocks.append(("head", _clean(st[4:])))
            continue
        if st.startswith("|") and st.endswith("|"):
            cells = [c.strip() for c in st.strip("|").split("|")]
            if all(set(c) <= set("-: ") for c in cells):
                continue
            tbl.append([_clean(c) for c in cells])
            continue
        flush_tbl()
        if not st or st == "---" or st.startswith("# ") or st.startswith("!["):
            continue
        if st.startswith("<") and st.endswith(">"):   # 순수 HTML 줄(제목·스타일 등)
            continue
        m = re.match(r"^(\s*)[-*] (.*)$", raw)
        if m:
            blocks.append(("bullet", 1 if len(m.group(1)) >= 2 else 0, _clean(m.group(2))))
            continue
        if st.startswith(">"):
            t = _clean(st.lstrip(">").strip())
            if t:
                blocks.append(("bullet", 0, t))
            continue
        blocks.append(("bullet", 0, _clean(st)))
    flush_tbl()
    if cur is not None:
        sections.append((cur, blocks))

    SIZE = {0: Pt(16), 1: Pt(13)}
    for title, blks in sections:
        pending = []

        def emit(cont):
            if not pending:
                return
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = title + (" (계속)" if cont else "")
            tf = sl.placeholders[1].text_frame
            tf.clear()
            first = True
            for b in pending:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                if b[0] == "head":
                    p.text, p.level = b[1], 0
                    p.font.bold, p.font.size = True, Pt(18)
                else:
                    p.text, p.level = b[2], b[1]
                    p.font.size = SIZE[b[1]]
            pending.clear()

        cont = False
        for b in blks:
            if b[0] == "table":
                emit(cont)
                cont = False
                _add_table_slide(prs, title, b[1])
                continue
            pending.append(b)
            if len(pending) >= 7:
                emit(cont)
                cont = True
        emit(cont)

    prs.save(str(ROOT / pptx_name))
    print(f"  PPTX {pptx_name}: {len(prs.slides)} slides")


def main():
    for md_name, slug, title, pdf, pptx in PAGES:
        body = convert(md_name)
        links = []
        if pdf:
            links.append(f'<a class="pdflink" href="{pdf}">⬇ PDF로 저장·인쇄</a>')
        if pptx:
            links.append(f'<a class="pdflink" href="{pptx}">⬇ PPTX 슬라이드</a>')
        if links:
            body = " ".join(links) + "\n" + body
        (ROOT / slug).write_text(TEMPLATE.format(title=title, css=CSS, body=body), encoding="utf-8")
        print("wrote", slug)
        if pdf:
            make_pdf(slug, pdf)
        if pptx:
            make_pptx(md_name, pptx, title)

    index_body = """<h1 align="center">모두가 주인공! 함께 성장하는 제주교육 AI 플랫폼</h1>
<p class="lead-sub">학생·교사·행정 모두 함께 쓰는 AI</p>
<p>미래학력분과 관련 자료 중 외부 공개용 참고 문서입니다. 아래에서 각 문서를 열람할 수 있습니다.</p>
<ul>
<li><a href="ax-transformation.html"><b>AX 전환과 교육</b></a><a class="idx-pdf" href="ax-transformation.pptx">PPTX</a><br>
AX(AI 전환)를 '도구 도입'이 아니라 사람·프로세스의 전환으로 보는 관점을 교육에 적용 — 교사 대체 우려, 에듀테크 과잉, 형평.</li>
<li><a href="paradigm.html"><b>AI 교육 패러다임</b></a><a class="idx-pdf" href="paradigm.pptx">PPTX</a><br>
AI 교육을 '무엇을'이 아니라 '어떻게'로 보는 관점 — 최고급 AI 접근의 공적 보장(토대), 활용 역량(핵심), 평가·교사 역할의 전환.</li>
<li><a href="bloom-2sigma.html"><b>고의숙 교육감의 "초개별화 맞춤형 교육"</b></a><a class="idx-pdf" href="bloom-2sigma.pdf">PDF</a><a class="idx-pdf" href="bloom-2sigma.pptx">PPTX</a><br>
핵심 공약 '초개별화 맞춤형 교육'의 이론적 근거 — 1:1 개인교습이 집단수업보다 2 표준편차 높은 성취를 낸다는 블룸의 실증(2 시그마 문제)을, AI로 대규모 실현한다.</li>
</ul>"""
    index_html = TEMPLATE.format(title="제주교육 AI 플랫폼", css=CSS, body=index_body)
    (ROOT / "index.html").write_text(index_html, encoding="utf-8")
    print("wrote index.html")
    print("빌드·검증 완료.")


if __name__ == "__main__":
    main()
