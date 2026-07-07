# -*- coding: utf-8 -*-
"""발표용 PPTX 생성 엔진 (16:9, 디자인 테마).

- 슬라이드 타입: title / section / stat / bullets / two / table / quote / closing
- 핵심 문서(bloom)는 CURATED 스펙으로 큐레이션, 그 외는 .md 에서 자동 요약.
- build.py 가 deck.make_pptx(md_path, out_path, title) 로 호출.
"""
import re
import html as _html
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 팔레트 ────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x2A, 0x44)
NAVY2  = RGBColor(0x17, 0x3B, 0x5C)
ACCENT = RGBColor(0x1C, 0xC5, 0xB0)   # teal
AMBER  = RGBColor(0xF4, 0xA2, 0x3B)
LIGHT  = RGBColor(0xF3, 0xF6, 0xFA)
INK    = RGBColor(0x1B, 0x2A, 0x38)
MUTED  = RGBColor(0x6E, 0x7B, 0x8A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LTXT   = RGBColor(0xC7, 0xD5, 0xE2)   # navy 배경 위 보조 텍스트
FONT   = "맑은 고딕"

W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.85)
CW = Inches(13.333 - 1.7)


# ── 저수준 헬퍼 ───────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, color, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", font)
    return r


def _par(tf, first, align=PP_ALIGN.LEFT, before=0, after=8, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if line:
        p.line_spacing = line
    return p


def _clean(s):
    s = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", s)
    s = re.sub(r"\[([^\]]+)\]\(<?[^>)]*>?\)", r"\1", s)
    s = s.replace("**", "").replace("`", "")
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    s = re.sub(r"<[^>]+>", "", s)
    return _html.unescape(s).strip()


# ── 슬라이드 빌더 ─────────────────────────────────────────
class Deck:
    def __init__(self, doc_short):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.doc_short = doc_short
        self.page = 0

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _footer(self, slide):
        self.page += 1
        tf = _box(slide, ML, Inches(7.02), Inches(9), Inches(0.35))
        _run(_par(tf, True, after=0), self.doc_short, 10.5, MUTED)
        tf2 = _box(slide, Inches(11.6), Inches(7.02), Inches(0.9), Inches(0.35))
        _run(_par(tf2, True, align=PP_ALIGN.RIGHT, after=0), str(self.page), 10.5, MUTED)

    # 표지
    def title(self, title, subtitle, footer, kicker="공개 참고 문서", watermark=""):
        s = self._blank()
        _bg(s, NAVY)
        _rect(s, 0, 0, W, Inches(0.22), ACCENT)
        if watermark:
            wm = _box(s, Inches(8.7), Inches(2.9), Inches(4.4), Inches(4.3), MSO_ANCHOR.BOTTOM)
            _run(_par(wm, True, align=PP_ALIGN.RIGHT, after=0), watermark, 220, NAVY2, bold=True)
        k = _box(s, ML, Inches(2.15), CW, Inches(0.5))
        _run(_par(k, True, after=0), kicker, 15, ACCENT, bold=True)
        t = _box(s, ML, Inches(2.65), Inches(10.6), Inches(1.9))
        _run(_par(t, True, after=0, line=1.06), title, 46, WHITE, bold=True)
        sub = _box(s, ML, Inches(4.75), Inches(10.6), Inches(1.4))
        _run(_par(sub, True, after=0, line=1.25), subtitle, 20, LTXT)
        _rect(s, ML, Inches(6.45), Inches(2.2), Inches(0.05), ACCENT)
        f = _box(s, ML, Inches(6.62), CW, Inches(0.5))
        _run(_par(f, True, after=0), footer, 14, LTXT)

    # 섹션 구분
    def section(self, no, title):
        s = self._blank()
        _bg(s, NAVY)
        _rect(s, 0, Inches(3.0), Inches(2.2), Inches(0.08), ACCENT)
        n = _box(s, ML, Inches(1.5), CW, Inches(1.5))
        _run(_par(n, True, after=0), no, 84, ACCENT, bold=True)
        t = _box(s, ML, Inches(3.25), Inches(11), Inches(2))
        _run(_par(t, True, after=0, line=1.1), title, 40, WHITE, bold=True)

    # 큰 수치
    def stat(self, value, caption, sub):
        s = self._blank()
        _bg(s, NAVY)
        v = _box(s, ML, Inches(1.3), Inches(5.4), Inches(4.3), MSO_ANCHOR.MIDDLE)
        _run(_par(v, True, align=PP_ALIGN.CENTER, after=0), value, 150, ACCENT, bold=True)
        _rect(s, Inches(6.35), Inches(2.35), Inches(0.06), Inches(2.9), ACCENT)
        c = _box(s, Inches(6.75), Inches(2.15), Inches(5.7), Inches(3.4), MSO_ANCHOR.TOP)
        _run(_par(c, True, after=10, line=1.15), caption, 25, WHITE, bold=True)
        _run(_par(c, False, after=0, line=1.3), sub, 16, LTXT)

    # 글머리표
    def bullets(self, title, items, kicker=None):
        # items: list of (level, text)
        s = self._blank()
        _bg(s, WHITE)
        _rect(s, 0, 0, W, Inches(0.16), ACCENT)
        if kicker:
            kk = _box(s, ML, Inches(0.42), CW, Inches(0.4))
            _run(_par(kk, True, after=0), kicker, 13, ACCENT, bold=True)
        t = _box(s, ML, Inches(0.78), CW, Inches(1.0))
        _run(_par(t, True, after=0), title, 29, NAVY, bold=True)
        _rect(s, ML, Inches(1.62), Inches(2.1), Inches(0.055), ACCENT)
        body = _box(s, ML, Inches(2.0), CW, Inches(4.7))
        first = True
        for lvl, text in items:
            p = _par(body, first, before=0, after=11 if lvl == 0 else 6, line=1.16)
            first = False
            if lvl == 0:
                _run(p, "●  ", 13, ACCENT, bold=True)
                _run(p, text, 18.5, INK)
            else:
                p.level = 1
                _run(p, "–  ", 13, MUTED, bold=True)
                _run(p, text, 15.5, MUTED)
        self._footer(s)

    # 좌우 2단 (라벨+내용)
    def two(self, title, left, right, kicker=None):
        s = self._blank()
        _bg(s, WHITE)
        _rect(s, 0, 0, W, Inches(0.16), ACCENT)
        t = _box(s, ML, Inches(0.6), CW, Inches(1.0))
        _run(_par(t, True, after=0), title, 29, NAVY, bold=True)
        _rect(s, ML, Inches(1.44), Inches(2.1), Inches(0.055), ACCENT)
        for i, (head, lines) in enumerate((left, right)):
            x = ML if i == 0 else Inches(7.05)
            card = _rect(s, x, Inches(1.95), Inches(5.5), Inches(4.35), LIGHT)
            hb = _box(s, x + Inches(0.35), Inches(2.25), Inches(4.8), Inches(0.7))
            _run(_par(hb, True, after=0), head, 19, NAVY, bold=True)
            bb = _box(s, x + Inches(0.35), Inches(3.05), Inches(4.85), Inches(3.1))
            fr = True
            for ln in lines:
                p = _par(bb, fr, after=9, line=1.18)
                fr = False
                _run(p, "· ", 14, ACCENT, bold=True)
                _run(p, ln, 15.5, INK)
        self._footer(s)

    # 표
    def table(self, title, rows, kicker=None):
        s = self._blank()
        _bg(s, WHITE)
        _rect(s, 0, 0, W, Inches(0.16), ACCENT)
        t = _box(s, ML, Inches(0.6), CW, Inches(1.0))
        _run(_par(t, True, after=0), title, 29, NAVY, bold=True)
        _rect(s, ML, Inches(1.44), Inches(2.1), Inches(0.055), ACCENT)
        nrows, ncols = len(rows), max(len(r) for r in rows)
        gt = s.shapes.add_table(nrows, ncols, ML, Inches(1.95), CW, Inches(0.5 * nrows)).table
        gt.first_row = False
        gt.horz_banding = False
        for i, row in enumerate(rows):
            gt.rows[i].height = Inches(0.62 if i == 0 else 0.55)
            for j in range(ncols):
                c = gt.cell(i, j)
                c.text = row[j] if j < len(row) else ""
                c.margin_left = Inches(0.15)
                c.margin_right = Inches(0.12)
                c.margin_top = Inches(0.03)
                c.margin_bottom = Inches(0.03)
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
                c.fill.solid()
                c.fill.fore_color.rgb = NAVY if i == 0 else (LIGHT if i % 2 == 0 else WHITE)
                for p in c.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    if not p.runs:
                        continue
                    for r in p.runs:
                        r.font.name = FONT
                        r.font.size = Pt(13 if i == 0 else 12.5)
                        r.font.bold = (i == 0)
                        r.font.color.rgb = WHITE if i == 0 else INK
                        rPr = r._r.get_or_add_rPr()
                        for tag in ("a:ea", "a:cs"):
                            e = rPr.find(qn(tag))
                            if e is None:
                                e = rPr.makeelement(qn(tag), {})
                                rPr.append(e)
                            e.set("typeface", FONT)
        self._footer(s)

    # 인용
    def quote(self, quote, attrib):
        s = self._blank()
        _bg(s, LIGHT)
        _rect(s, 0, 0, Inches(0.22), H, ACCENT)
        qm = _box(s, ML, Inches(0.7), Inches(3), Inches(1.6))
        _run(_par(qm, True, after=0), "“", 130, ACCENT, bold=True)
        q = _box(s, Inches(1.4), Inches(2.1), Inches(10.6), Inches(3.4), MSO_ANCHOR.TOP)
        _run(_par(q, True, after=0, line=1.35), quote, 27, INK, bold=True)
        a = _box(s, Inches(1.4), Inches(5.6), Inches(10.6), Inches(0.8))
        _run(_par(a, True, after=0), attrib, 17, MUTED)

    # 맺음
    def closing(self, text, footer):
        s = self._blank()
        _bg(s, NAVY)
        _rect(s, 0, 0, W, Inches(0.22), ACCENT)
        k = _box(s, ML, Inches(2.0), CW, Inches(0.5))
        _run(_par(k, True, after=0), "한 줄로", 16, ACCENT, bold=True)
        t = _box(s, ML, Inches(2.6), Inches(11.4), Inches(3.2))
        _run(_par(t, True, after=0, line=1.32), text, 27, WHITE, bold=True)
        _rect(s, ML, Inches(6.35), Inches(2.2), Inches(0.05), ACCENT)
        f = _box(s, ML, Inches(6.55), CW, Inches(0.5))
        _run(_par(f, True, after=0), footer, 14, LTXT)

    def save(self, path):
        self.prs.save(path)
        return len(self.prs.slides)


# ── 자동 요약(비큐레이션 문서) ────────────────────────────
def _sections(md_text):
    lines = md_text.split("\n")
    sections, cur, blocks, tbl, code = [], None, [], [], False
    def flush():
        if tbl:
            blocks.append(("table", list(tbl)))
            tbl.clear()
    for raw in lines:
        st = raw.strip()
        if st.startswith("```"):
            code = not code; continue
        if code:
            continue
        if st.startswith("## "):
            flush()
            if cur is not None:
                sections.append((cur, blocks))
            cur, blocks = _clean(st[3:]), []
            continue
        if st.startswith("### "):
            flush(); blocks.append(("head", _clean(st[4:]))); continue
        if st.startswith("|") and st.endswith("|"):
            cells = [c.strip() for c in st.strip("|").split("|")]
            if all(set(c) <= set("-: ") for c in cells):
                continue
            tbl.append([_clean(c) for c in cells]); continue
        flush()
        if not st or st == "---" or st.startswith("# ") or st.startswith("!["):
            continue
        if st.startswith("<") and st.endswith(">"):
            continue
        m = re.match(r"^(\s*)[-*] (.*)$", raw)
        if m:
            blocks.append(("bullet", 1 if len(m.group(1)) >= 2 else 0, _clean(m.group(2)))); continue
        if st.startswith(">"):
            t = _clean(st.lstrip(">").strip())
            if t:
                blocks.append(("bullet", 0, t))
            continue
        blocks.append(("bullet", 0, _clean(st)))
    flush()
    if cur is not None:
        sections.append((cur, blocks))
    return sections


def _short(t, n=150):
    if len(t) <= n:
        return t
    cut = t[:n]
    for sep in ("다. ", ". ", "—", "; ", ", "):
        i = cut.rfind(sep)
        if i > 70:
            end = i + (2 if sep == "다. " or sep == ". " else 0)
            return cut[:end].strip().rstrip(",") + "…"
    return cut.strip() + "…"


def _auto(md_text, deck_title, doc_short):
    d = Deck(doc_short)
    d.title(deck_title, "제주교육 AI · 미래학력분과 참고 자료",
            "제주도교육인수위원회 · 미래학력분과")
    for title, blocks in _sections(md_text):
        pend, first_slide = [], True
        def emit(cont):
            nonlocal pend
            if not pend:
                return
            d.bullets(title + ("" if not cont else " (계속)"), pend)
            pend = []
        for b in blocks:
            if b[0] == "table":
                emit(not first_slide)
                d.table(title, b[1])
                first_slide = False
                continue
            if b[0] == "head":
                pend.append((0, b[1]))
            else:
                lvl, txt = b[1], _short(b[2])
                pend.append((lvl, txt))
            if len(pend) >= 5:
                emit(not first_slide)
                first_slide = False
        emit(not first_slide)
    return d


# ── 큐레이션 스펙 ─────────────────────────────────────────
def _bloom(_deck_title):
    d = Deck("블룸 2σ · 초개별화 맞춤형 교육")
    d.title(
        "초개별화 맞춤형 교육",
        "블룸의 ‘2 시그마 문제’와 AI — 고의숙 교육감 핵심공약의 이론적 배경",
        "제주도교육인수위원회 · 미래학력분과",
        kicker="이론적 배경 · 핵심공약",
        watermark="2σ",
    )
    d.section("01", "문제 — ‘2 시그마’")
    d.stat(
        "2σ",
        "1:1 개인교습을 받은 학생은 집단수업보다 약 2 표준편차 높은 성취를 낸다",
        "그러나 개인교습을 모든 학생에게 제공하기엔 비용이 너무 크다 — 블룸(1984)이 던진 ‘2 시그마 문제’.",
    )
    d.two(
        "실험 설계 — 세 조건 무작위 비교",
        ("설계", [
            "시카고대 아나니아·버크(1982~84)",
            "4·5·8학년, 3주·11차시 반복 검증",
            "초기 적성·수업시간 동일하게 통제",
        ]),
        ("세 가지 학습 조건", [
            "① 전통 집단수업 (교사 1 : 다수)",
            "② 완전학습 (형성평가 + 교정)",
            "③ 1:1·소집단 개인교습",
        ]),
    )
    d.table(
        "결과 — ‘2 시그마’ 효과",
        [
            ["학습 조건", "성취(전통 대비)", "목표 도달률"],
            ["전통 집단수업", "기준", "상위 20%"],
            ["완전학습", "+1σ", "70%"],
            ["1:1 개인교습", "+2σ", "90%"],
        ],
    )
    d.bullets(
        "무엇을 뜻하는가",
        [
            (0, "타고난 적성보다 ‘학습 조건’이 성취를 결정한다 (적성–성취 상관 +.60 → +.25)"),
            (0, "대부분의 학생이 높은 수준에 도달할 잠재력을 가지고 있다 — 블룸"),
            (0, "핵심 과제: 1:1 교습에 근접하면서도 대규모로 가능한 방법을 찾는 것"),
        ],
        kicker="핵심 메시지",
    )
    d.quote(
        "일대일 개인교습에 근접한 효과를 내면서도, 대규모로 실현 가능한 집단 교수·학습 조건을 찾을 수 있는가?",
        "— Benjamin Bloom, ‘The 2 Sigma Problem’ (1984)",
    )
    d.section("02", "해법 — AI, 고전에서 실증으로")
    d.bullets(
        "① 생성형 AI 이전 — 지능형 튜터링(ITS)",
        [
            (0, "VanLehn(2011): 단계형 ITS는 인간 교사에 근접한 효과"),
            (0, "Kulik & Fletcher(2016): ITS 중앙값 +0.66σ, 92% 연구에서 전통수업 능가"),
            (1, "‘2σ’는 인간만의 값이 아니라, 잘 설계된 기술로 좁힐 수 있는 격차"),
        ],
        kicker="현대적 실증",
    )
    d.bullets(
        "② 생성형 AI 시대(2023~) — 무작위통제연구(RCT)",
        [
            (0, "Kestin 외(2025, 하버드): 잘 설계된 AI 튜터가 활동중심 수업보다 2배 이상 학습 이득 — 더 짧은 시간에"),
            (0, "나이지리아 대규모 RCT(World Bank, 2024): 생성형 AI 튜터 영어수업, 유의한 성취 향상"),
            (1, "공통점 — “AI를 쓰면”이 아니라 “잘 설계된 AI가” 오른다"),
        ],
        kicker="현대적 실증",
    )
    d.bullets(
        "어떻게 2σ를 푸는가 — 작동 원리",
        [
            (0, "생성형 AI = 자연어 대화 + 실시간 적응 + 낮은 한계비용 → 블룸의 ‘비용 장벽’을 낮춘다"),
            (0, "AI는 학생마다 숙달을 향해 「형성평가 → 교정」 루프를 도는 개인 튜터"),
            (0, "교사는 서로 다른 속도의 여러 학생을 조율하는 지휘자"),
            (0, "초개별화 = 완전학습의 개별화 + AI 적응형 대화를, 대규모로"),
        ],
    )
    d.two(
        "그러나 신중하게 — ‘콘텐츠’가 아니라 ‘조건’",
        ("무엇에 투자하나", [
            "ⓐ 좋은 재료 — 고성능 AI의 공적 보장",
            "ⓑ 좋은 주방 — 도구·전량 기록·환류",
            "ⓒ 교사 역량 — 좋은 수업 설계",
        ]),
        ("무엇이 아닌가", [
            "제주가 콘텐츠를 직접 찍어내는 일이 아니다",
            "콘텐츠는 교사·확장 생태계의 몫",
            "근거는 축적 중 → 자체 시범으로 검증",
        ]),
    )
    d.section("03", "제주 초개별화로")
    d.bullets(
        "시사점 — 핵심공약과의 연결",
        [
            (0, "AI가 ‘2 시그마 문제’의 현실적 해법 — 공약의 논리적 출발점"),
            (0, "초개별화의 이론적 뿌리 = 블룸 완전학습(형성평가+교정)의 AI·디지털 구현"),
            (0, "형평성: “적성보다 학습 조건” → AI 접근 공적 보장이 출발선·격차 해소"),
            (0, "목표: “대부분의 학생이 높은 수준 도달 가능” = ‘한 아이도 놓치지 않는’ 책임교육"),
        ],
    )
    d.closing(
        "블룸(1984)이 비용 때문에 미룬 ‘1:1 교습의 대규모 실현’을, AI가 실시간 개별화와 환류로 푼다 — 학생 한 명 한 명에게 개인교습에 준하는 조건을, 대규모로.",
        "제주도교육인수위원회 · 미래학력분과",
    )
    return d


CURATED = {"bloom-2sigma.pptx": _bloom}


# ── 진입점 ────────────────────────────────────────────────
def make_pptx(md_path, out_path, deck_title):
    import os
    name = os.path.basename(out_path)
    short = re.sub(r"\.pptx$", "", name)
    if name in CURATED:
        d = CURATED[name](deck_title)
    else:
        md_text = open(md_path, encoding="utf-8").read()
        d = _auto(md_text, _clean(deck_title), deck_title.split("—")[0].strip() if "—" in deck_title else deck_title)
    n = d.save(out_path)
    print(f"  PPTX {name}: {n} slides")
